#!/usr/bin/env python3
"""Tests for the standalone PowerPoint slide-transfer tool."""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches


_MINIMAL_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cc\xf8\xff\xff?"
    b"\x00\x05\xfe\x02\xfe\xa3\x9b\x14\xa1\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _count_package_members(path: Path, prefix: str) -> int:
    with zipfile.ZipFile(path, "r") as zf:
        return sum(1 for name in zf.namelist() if name.startswith(prefix))


def _build_source_with_picture(path: Path, image_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Imported Title"
    body = slide.placeholders[1].text_frame
    body.text = "Imported body"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), Inches(1.5), Inches(1.5))
    prs.save(path)


def _build_target(path: Path, titles: list[str]) -> None:
    prs = Presentation()
    for title in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Body for {title}"
    prs.save(path)


def test_import_slide_copies_slide_assets_and_reuses_default_layout_master(
    pptx_slide_transfer_tools,
    temp_dir,
):
    image_path = temp_dir / "tiny.png"
    image_path.write_bytes(_MINIMAL_PNG)

    source = temp_dir / "source.pptx"
    target = temp_dir / "target.pptx"
    _build_source_with_picture(source, image_path)
    _build_target(target, ["Target One"])

    masters_before = _count_package_members(target, "ppt/slideMasters/slideMaster")
    media_before = _count_package_members(target, "ppt/media/")

    result = pptx_slide_transfer_tools.tool_pptx_import_slide(
        str(source),
        1,
        str(target),
    )
    assert result["success"] is True
    assert result["new_slide_number"] == 2
    assert result["master_copied"] is False
    assert result["layout_reused"] is True

    prs = Presentation(str(target))
    assert len(prs.slides) == 2
    imported = prs.slides[1]
    assert imported.shapes.title.text == "Imported Title"
    assert any(getattr(shape, "shape_type", None) == 13 for shape in imported.shapes)  # picture

    masters_after = _count_package_members(target, "ppt/slideMasters/slideMaster")
    media_after = _count_package_members(target, "ppt/media/")
    assert masters_after == masters_before
    assert media_after >= media_before + 1


def test_import_slide_after_specific_position_preserves_order(
    pptx_slide_transfer_tools,
    temp_dir,
):
    image_path = temp_dir / "tiny.png"
    image_path.write_bytes(_MINIMAL_PNG)

    source = temp_dir / "source_order.pptx"
    target = temp_dir / "target_order.pptx"
    _build_source_with_picture(source, image_path)
    _build_target(target, ["First", "Second"])

    result = pptx_slide_transfer_tools.tool_pptx_import_slide(
        str(source),
        1,
        str(target),
        position="after",
        after_slide_number=1,
    )
    assert result["success"] is True
    assert result["new_slide_number"] == 2

    prs = Presentation(str(target))
    titles = [slide.shapes.title.text for slide in prs.slides]
    assert titles == ["First", "Imported Title", "Second"]


def test_importing_same_source_slide_twice_does_not_duplicate_default_master(
    pptx_slide_transfer_tools,
    temp_dir,
):
    image_path = temp_dir / "tiny.png"
    image_path.write_bytes(_MINIMAL_PNG)

    source = temp_dir / "source_twice.pptx"
    target = temp_dir / "target_twice.pptx"
    _build_source_with_picture(source, image_path)
    _build_target(target, ["Seed"])

    result1 = pptx_slide_transfer_tools.tool_pptx_import_slide(str(source), 1, str(target))
    result2 = pptx_slide_transfer_tools.tool_pptx_import_slide(str(source), 1, str(target))
    assert result1["success"] is True
    assert result2["success"] is True

    prs = Presentation(str(target))
    assert len(prs.slides) == 3
    assert _count_package_members(target, "ppt/slideMasters/slideMaster") == 1


def test_import_slide_rejects_invalid_source_slide_number(
    pptx_slide_transfer_tools,
    temp_dir,
):
    source = temp_dir / "source_invalid.pptx"
    target = temp_dir / "target_invalid.pptx"
    _build_target(source, ["Only"])
    _build_target(target, ["Target"])

    result = pptx_slide_transfer_tools.tool_pptx_import_slide(
        str(source),
        2,
        str(target),
    )
    assert "error" in result
    assert "Presentation has 1 slides" in result["error"]


def test_import_slide_requires_after_slide_number_for_after_mode(
    pptx_slide_transfer_tools,
    temp_dir,
):
    source = temp_dir / "source_after.pptx"
    target = temp_dir / "target_after.pptx"
    _build_target(source, ["Source"])
    _build_target(target, ["Target"])

    result = pptx_slide_transfer_tools.tool_pptx_import_slide(
        str(source),
        1,
        str(target),
        position="after",
    )
    assert result == {"error": "after_slide_number is required when position='after'."}
