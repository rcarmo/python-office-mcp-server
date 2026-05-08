#!/usr/bin/env python3
"""Standalone PowerPoint slide transfer tool.

Copies a single slide from one presentation to another together with the
parts it depends on (images, charts, embedded packages, media, layout,
master, theme, and related OOXML parts needed for rendering).

Design notes:

* If the source slide's *layout* already exists verbatim in the target
  package, reuse it and do **not** import a new layout/master/theme.
* Otherwise, import exactly one layout and exactly one master -- the one
  required by the slide. The copied master is pruned so it only keeps
  the imported layout relationship instead of dragging every layout from
  the source presentation across.
* Non-render dependencies such as notes slides are optional and excluded
  by default.

This stays at the OOXML package level rather than using python-pptx's
object model for the copy itself, because cross-package slide copying is
not something python-pptx exposes directly.
"""

from __future__ import annotations

import posixpath
import re
import shutil
import tempfile
import zipfile
from copy import deepcopy
from pathlib import Path
from typing import Any

try:
    from lxml import etree
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .save_utils import resolve_office_path

PKG_CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DOC_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NSMAP = {
    "ct": PKG_CT_NS,
    "rel": PKG_REL_NS,
    "p": PML_NS,
    "r": DOC_REL_NS,
}

RT_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
RT_SLIDE_LAYOUT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
RT_SLIDE_MASTER = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"
RT_THEME = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
RT_NOTES_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"


def _read_package(path: str) -> dict[str, bytes]:
    with zipfile.ZipFile(path, "r") as zf:
        return {name: zf.read(name) for name in zf.namelist()}


def _write_package(parts: dict[str, bytes], output_path: str) -> None:
    fd, tmp_path = tempfile.mkstemp(suffix=Path(output_path).suffix or ".pptx", dir=Path(output_path).parent)
    Path(tmp_path).unlink(missing_ok=True)
    try:
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for name in sorted(parts):
                zf.writestr(name, parts[name])
        shutil.move(tmp_path, output_path)
    finally:
        if Path(tmp_path).exists():
            Path(tmp_path).unlink(missing_ok=True)


def _rels_path_for(part_path: str) -> str:
    parent = posixpath.dirname(part_path)
    name = posixpath.basename(part_path)
    if parent:
        return posixpath.join(parent, "_rels", name + ".rels")
    return posixpath.join("_rels", name + ".rels")


def _resolve_rel_target(base_part: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    base_dir = posixpath.dirname(base_part)
    return posixpath.normpath(posixpath.join(base_dir, target)).lstrip("/")


def _relative_target(from_part: str, to_part: str) -> str:
    base_dir = posixpath.dirname(from_part) or "."
    return posixpath.relpath(to_part, base_dir)


def _parse_xml(data: bytes):
    return etree.fromstring(data)


def _serialise_xml(root) -> bytes:
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone="yes")


def _find_override_content_type(ct_root, part_path: str) -> str | None:
    wanted = "/" + part_path
    for node in ct_root.findall("ct:Override", namespaces=NSMAP):
        if node.get("PartName") == wanted:
            return node.get("ContentType")
    return None


def _find_default_content_type(ct_root, ext: str) -> str | None:
    for node in ct_root.findall("ct:Default", namespaces=NSMAP):
        if node.get("Extension", "").lower() == ext.lower().lstrip("."):
            return node.get("ContentType")
    return None


def _ensure_content_type(src_ct_root, dst_ct_root, src_part: str, dst_part: str) -> None:
    override = _find_override_content_type(src_ct_root, src_part)
    if override:
        wanted = "/" + dst_part
        existing = [
            node for node in dst_ct_root.findall("ct:Override", namespaces=NSMAP)
            if node.get("PartName") == wanted
        ]
        if not existing:
            node = etree.SubElement(dst_ct_root, f"{{{PKG_CT_NS}}}Override")
            node.set("PartName", wanted)
            node.set("ContentType", override)
        return

    ext = Path(dst_part).suffix.lstrip(".")
    if not ext:
        return
    default_ct = _find_default_content_type(src_ct_root, ext)
    if not default_ct:
        return
    existing = [
        node for node in dst_ct_root.findall("ct:Default", namespaces=NSMAP)
        if node.get("Extension", "").lower() == ext.lower()
    ]
    if not existing:
        node = etree.SubElement(dst_ct_root, f"{{{PKG_CT_NS}}}Default")
        node.set("Extension", ext)
        node.set("ContentType", default_ct)


def _next_part_name(parts: dict[str, bytes], directory: str, stem: str, suffix: str) -> str:
    stem = re.sub(r"\d+$", "", stem) or stem
    idx = 1
    while True:
        candidate = posixpath.join(directory, f"{stem}{idx}{suffix}")
        if candidate not in parts:
            return candidate
        idx += 1


def _clone_rels_root(parts: dict[str, bytes], src_part: str):
    rels_path = _rels_path_for(src_part)
    if rels_path not in parts:
        return None
    return _parse_xml(parts[rels_path])


def _find_relationship(root, reltype: str):
    if root is None:
        return None
    for rel in root.findall("rel:Relationship", namespaces=NSMAP):
        if rel.get("Type") == reltype:
            return rel
    return None


def _find_matching_layout_by_bytes(dst_parts: dict[str, bytes], layout_bytes: bytes) -> str | None:
    for part_name, part_bytes in dst_parts.items():
        if part_name.startswith("ppt/slideLayouts/slideLayout") and part_name.endswith(".xml"):
            if part_bytes == layout_bytes:
                return part_name
    return None


def _copy_generic_part(
    src_parts: dict[str, bytes],
    dst_parts: dict[str, bytes],
    src_ct_root,
    dst_ct_root,
    src_part: str,
    mapping: dict[str, str],
    *,
    skip_reltypes: set[str] | None = None,
) -> str:
    if src_part in mapping:
        return mapping[src_part]

    directory = posixpath.dirname(src_part)
    stem = Path(src_part).stem
    suffix = Path(src_part).suffix
    dst_part = _next_part_name(dst_parts, directory, stem, suffix)
    mapping[src_part] = dst_part
    dst_parts[dst_part] = src_parts[src_part]
    _ensure_content_type(src_ct_root, dst_ct_root, src_part, dst_part)

    rels_root = _clone_rels_root(src_parts, src_part)
    if rels_root is None:
        return dst_part

    skip_reltypes = skip_reltypes or set()
    for rel in list(rels_root.findall("rel:Relationship", namespaces=NSMAP)):
        if rel.get("Type") in skip_reltypes:
            rel.getparent().remove(rel)
            continue
        if rel.get("TargetMode") == "External":
            continue
        child_src = _resolve_rel_target(src_part, rel.get("Target"))
        child_dst = _copy_generic_part(src_parts, dst_parts, src_ct_root, dst_ct_root, child_src, mapping)
        rel.set("Target", _relative_target(dst_part, child_dst))

    dst_parts[_rels_path_for(dst_part)] = _serialise_xml(rels_root)
    return dst_part


def _copy_theme_part(
    src_parts: dict[str, bytes],
    dst_parts: dict[str, bytes],
    src_ct_root,
    dst_ct_root,
    src_theme_part: str,
    mapping: dict[str, str],
) -> str:
    if src_theme_part in mapping:
        return mapping[src_theme_part]

    for part_name, part_bytes in dst_parts.items():
        if part_name.startswith("ppt/theme/theme") and part_name.endswith(".xml"):
            if part_bytes == src_parts[src_theme_part]:
                mapping[src_theme_part] = part_name
                return part_name

    return _copy_generic_part(src_parts, dst_parts, src_ct_root, dst_ct_root, src_theme_part, mapping)


def _copy_master_and_layout(
    src_parts: dict[str, bytes],
    dst_parts: dict[str, bytes],
    src_ct_root,
    dst_ct_root,
    src_layout_part: str,
    mapping: dict[str, str],
) -> tuple[str, str, bool]:
    """Copy the layout+master chain, or reuse an identical layout if present.

    Returns: (target_layout_part, target_master_part, master_was_copied)
    """
    reused_layout = _find_matching_layout_by_bytes(dst_parts, src_parts[src_layout_part])
    if reused_layout:
        layout_rels = _clone_rels_root(dst_parts, reused_layout)
        master_rel = _find_relationship(layout_rels, RT_SLIDE_MASTER)
        if master_rel is None:
            raise ValueError(f"Reused layout {reused_layout} has no slideMaster relationship")
        master_part = _resolve_rel_target(reused_layout, master_rel.get("Target"))
        mapping[src_layout_part] = reused_layout
        return reused_layout, master_part, False

    src_layout_rels = _clone_rels_root(src_parts, src_layout_part)
    if src_layout_rels is None:
        raise ValueError(f"Source layout {src_layout_part} has no relationships part")
    src_master_rel = _find_relationship(src_layout_rels, RT_SLIDE_MASTER)
    if src_master_rel is None:
        raise ValueError(f"Source layout {src_layout_part} has no slideMaster relationship")
    src_master_part = _resolve_rel_target(src_layout_part, src_master_rel.get("Target"))

    src_master_rels = _clone_rels_root(src_parts, src_master_part)
    if src_master_rels is None:
        raise ValueError(f"Source master {src_master_part} has no relationships part")

    src_master_layout_rel = None
    for rel in src_master_rels.findall("rel:Relationship", namespaces=NSMAP):
        if rel.get("Type") != RT_SLIDE_LAYOUT or rel.get("TargetMode") == "External":
            continue
        if _resolve_rel_target(src_master_part, rel.get("Target")) == src_layout_part:
            src_master_layout_rel = rel
            break
    if src_master_layout_rel is None:
        raise ValueError(
            f"Source master {src_master_part} does not reference layout {src_layout_part}"
        )

    # Copy / reuse the theme first so the master can point to it.
    target_theme_part: str | None = None
    src_theme_rel = _find_relationship(src_master_rels, RT_THEME)
    if src_theme_rel is not None and src_theme_rel.get("TargetMode") != "External":
        src_theme_part = _resolve_rel_target(src_master_part, src_theme_rel.get("Target"))
        target_theme_part = _copy_theme_part(src_parts, dst_parts, src_ct_root, dst_ct_root, src_theme_part, mapping)

    target_master_part = _next_part_name(dst_parts, "ppt/slideMasters", "slideMaster", ".xml")
    target_layout_part = _next_part_name(dst_parts, "ppt/slideLayouts", "slideLayout", ".xml")
    mapping[src_master_part] = target_master_part
    mapping[src_layout_part] = target_layout_part

    # Copy the layout and rewrite its master relationship plus any non-master deps.
    dst_parts[target_layout_part] = src_parts[src_layout_part]
    _ensure_content_type(src_ct_root, dst_ct_root, src_layout_part, target_layout_part)
    layout_rels = deepcopy(src_layout_rels)
    for rel in layout_rels.findall("rel:Relationship", namespaces=NSMAP):
        if rel.get("Type") == RT_SLIDE_MASTER:
            rel.set("Target", _relative_target(target_layout_part, target_master_part))
            continue
        if rel.get("TargetMode") == "External":
            continue
        child_src = _resolve_rel_target(src_layout_part, rel.get("Target"))
        child_dst = _copy_generic_part(src_parts, dst_parts, src_ct_root, dst_ct_root, child_src, mapping)
        rel.set("Target", _relative_target(target_layout_part, child_dst))
    dst_parts[_rels_path_for(target_layout_part)] = _serialise_xml(layout_rels)

    # Copy the master, but keep only the imported layout relationship.
    master_xml = _parse_xml(src_parts[src_master_part])
    layout_id_lst = master_xml.find("p:sldLayoutIdLst", namespaces=NSMAP)
    if layout_id_lst is not None:
        for child in list(layout_id_lst):
            if child.get(f"{{{DOC_REL_NS}}}id") != src_master_layout_rel.get("Id"):
                layout_id_lst.remove(child)

    dst_parts[target_master_part] = _serialise_xml(master_xml)
    _ensure_content_type(src_ct_root, dst_ct_root, src_master_part, target_master_part)

    master_rels = deepcopy(src_master_rels)
    for rel in list(master_rels.findall("rel:Relationship", namespaces=NSMAP)):
        if rel.get("Type") == RT_SLIDE_LAYOUT:
            if rel.get("Id") == src_master_layout_rel.get("Id"):
                rel.set("Target", _relative_target(target_master_part, target_layout_part))
            else:
                rel.getparent().remove(rel)
            continue
        if rel.get("Type") == RT_THEME:
            if target_theme_part is not None:
                rel.set("Target", _relative_target(target_master_part, target_theme_part))
            else:
                rel.getparent().remove(rel)
            continue
        if rel.get("TargetMode") == "External":
            continue
        child_src = _resolve_rel_target(src_master_part, rel.get("Target"))
        child_dst = _copy_generic_part(src_parts, dst_parts, src_ct_root, dst_ct_root, child_src, mapping)
        rel.set("Target", _relative_target(target_master_part, child_dst))
    dst_parts[_rels_path_for(target_master_part)] = _serialise_xml(master_rels)

    return target_layout_part, target_master_part, True


def _copy_slide_part(
    src_parts: dict[str, bytes],
    dst_parts: dict[str, bytes],
    src_ct_root,
    dst_ct_root,
    src_slide_part: str,
    target_layout_part: str,
    mapping: dict[str, str],
    *,
    include_notes: bool,
) -> str:
    target_slide_part = _next_part_name(dst_parts, "ppt/slides", "slide", ".xml")
    mapping[src_slide_part] = target_slide_part
    dst_parts[target_slide_part] = src_parts[src_slide_part]
    _ensure_content_type(src_ct_root, dst_ct_root, src_slide_part, target_slide_part)

    slide_rels = _clone_rels_root(src_parts, src_slide_part)
    if slide_rels is None:
        return target_slide_part

    for rel in list(slide_rels.findall("rel:Relationship", namespaces=NSMAP)):
        if rel.get("Type") == RT_SLIDE_LAYOUT:
            rel.set("Target", _relative_target(target_slide_part, target_layout_part))
            continue
        if rel.get("Type") == RT_NOTES_SLIDE and not include_notes:
            rel.getparent().remove(rel)
            continue
        if rel.get("TargetMode") == "External":
            continue
        child_src = _resolve_rel_target(src_slide_part, rel.get("Target"))
        child_dst = _copy_generic_part(src_parts, dst_parts, src_ct_root, dst_ct_root, child_src, mapping)
        rel.set("Target", _relative_target(target_slide_part, child_dst))

    dst_parts[_rels_path_for(target_slide_part)] = _serialise_xml(slide_rels)
    return target_slide_part


def _next_rel_id(rels_root) -> str:
    ids = []
    for rel in rels_root.findall("rel:Relationship", namespaces=NSMAP):
        rid = rel.get("Id", "")
        m = re.match(r"rId(\d+)$", rid)
        if m:
            ids.append(int(m.group(1)))
    return f"rId{max(ids, default=0) + 1}"


def _next_slide_id(presentation_root) -> int:
    ids = []
    for node in presentation_root.findall("p:sldIdLst/p:sldId", namespaces=NSMAP):
        try:
            ids.append(int(node.get("id")))
        except (TypeError, ValueError):
            continue
    return max(ids, default=255) + 1


def _next_master_id(presentation_root) -> int:
    ids = []
    for node in presentation_root.findall("p:sldMasterIdLst/p:sldMasterId", namespaces=NSMAP):
        try:
            ids.append(int(node.get("id")))
        except (TypeError, ValueError):
            continue
    return max(ids, default=2147483647) + 1


def _find_slide_part_by_number(parts: dict[str, bytes], slide_number: int) -> tuple[str, int]:
    presentation_root = _parse_xml(parts["ppt/presentation.xml"])
    presentation_rels = _parse_xml(parts["ppt/_rels/presentation.xml.rels"])
    slide_ids = presentation_root.findall("p:sldIdLst/p:sldId", namespaces=NSMAP)
    if slide_number < 1 or slide_number > len(slide_ids):
        raise IndexError(f"Slide {slide_number} not found. Presentation has {len(slide_ids)} slides.")

    slide_id_node = slide_ids[slide_number - 1]
    rid = slide_id_node.get(f"{{{DOC_REL_NS}}}id")
    for rel in presentation_rels.findall("rel:Relationship", namespaces=NSMAP):
        if rel.get("Id") == rid:
            return _resolve_rel_target("ppt/presentation.xml", rel.get("Target")), len(slide_ids)
    raise ValueError(f"Could not resolve slide relationship {rid}")


def _find_slide_layout_part(src_parts: dict[str, bytes], slide_part: str) -> str:
    rels_root = _clone_rels_root(src_parts, slide_part)
    if rels_root is None:
        raise ValueError(f"Slide part {slide_part} has no relationships")
    rel = _find_relationship(rels_root, RT_SLIDE_LAYOUT)
    if rel is None:
        raise ValueError(f"Slide part {slide_part} has no slideLayout relationship")
    return _resolve_rel_target(slide_part, rel.get("Target"))


def _register_master_if_needed(dst_parts: dict[str, bytes], target_master_part: str) -> bool:
    presentation_root = _parse_xml(dst_parts["ppt/presentation.xml"])
    presentation_rels = _parse_xml(dst_parts["ppt/_rels/presentation.xml.rels"])

    for rel in presentation_rels.findall("rel:Relationship", namespaces=NSMAP):
        if rel.get("Type") == RT_SLIDE_MASTER:
            existing = _resolve_rel_target("ppt/presentation.xml", rel.get("Target"))
            if existing == target_master_part:
                return False

    rel_id = _next_rel_id(presentation_rels)
    rel = etree.SubElement(presentation_rels, f"{{{PKG_REL_NS}}}Relationship")
    rel.set("Id", rel_id)
    rel.set("Type", RT_SLIDE_MASTER)
    rel.set("Target", _relative_target("ppt/presentation.xml", target_master_part))

    master_list = presentation_root.find("p:sldMasterIdLst", namespaces=NSMAP)
    if master_list is None:
        master_list = etree.SubElement(presentation_root, f"{{{PML_NS}}}sldMasterIdLst")
    master_node = etree.SubElement(master_list, f"{{{PML_NS}}}sldMasterId")
    master_node.set("id", str(_next_master_id(presentation_root)))
    master_node.set(f"{{{DOC_REL_NS}}}id", rel_id)

    dst_parts["ppt/presentation.xml"] = _serialise_xml(presentation_root)
    dst_parts["ppt/_rels/presentation.xml.rels"] = _serialise_xml(presentation_rels)
    return True


def _insert_slide_into_presentation(
    dst_parts: dict[str, bytes],
    target_slide_part: str,
    *,
    position: str,
    after_slide_number: int | None,
) -> int:
    presentation_root = _parse_xml(dst_parts["ppt/presentation.xml"])
    presentation_rels = _parse_xml(dst_parts["ppt/_rels/presentation.xml.rels"])

    rel_id = _next_rel_id(presentation_rels)
    rel = etree.SubElement(presentation_rels, f"{{{PKG_REL_NS}}}Relationship")
    rel.set("Id", rel_id)
    rel.set("Type", RT_SLIDE)
    rel.set("Target", _relative_target("ppt/presentation.xml", target_slide_part))

    slide_list = presentation_root.find("p:sldIdLst", namespaces=NSMAP)
    if slide_list is None:
        slide_list = etree.SubElement(presentation_root, f"{{{PML_NS}}}sldIdLst")

    slide_node = etree.Element(f"{{{PML_NS}}}sldId")
    slide_node.set("id", str(_next_slide_id(presentation_root)))
    slide_node.set(f"{{{DOC_REL_NS}}}id", rel_id)

    if position == "after":
        if after_slide_number is None:
            raise ValueError("after_slide_number is required when position='after'")
        children = list(slide_list)
        if after_slide_number < 1 or after_slide_number > len(children):
            raise IndexError(
                f"after_slide_number {after_slide_number} invalid. Presentation has {len(children)} slides."
            )
        slide_list.insert(after_slide_number, slide_node)
        new_slide_number = after_slide_number + 1
    else:
        slide_list.append(slide_node)
        new_slide_number = len(list(slide_list))

    dst_parts["ppt/presentation.xml"] = _serialise_xml(presentation_root)
    dst_parts["ppt/_rels/presentation.xml.rels"] = _serialise_xml(presentation_rels)
    return new_slide_number


class PresentationSlideTransferTools:
    """Standalone tools for slide transfer across presentations."""

    def tool_pptx_import_slide(
        self,
        source_file_path: str,
        source_slide_number: int,
        target_file_path: str,
        position: str = "end",
        after_slide_number: int | None = None,
        output_path: str | None = None,
        include_notes: bool = False,
    ) -> dict[str, Any]:
        """Copy one slide from a source presentation into a target presentation.

        USE THIS when you need to lift a fully-designed slide across from one
        deck to another, preserving pictures, charts, embedded assets, and the
        layout/master/theme chain required for it to render correctly.

        The tool validates the source slide's layout/master chain before the
        copy. If an identical layout already exists in the target package, it
        reuses it and does not import a new master. Otherwise it imports only
        the single layout and single master required by the slide.

        Args:
            source_file_path: Source .pptx file containing the slide to copy
            source_slide_number: 1-based slide number in the source deck
            target_file_path: Target .pptx file to receive the slide
            position: 'end' or 'after'
            after_slide_number: Required when position='after'
            output_path: Optional output path (defaults to overwriting target)
            include_notes: Whether to carry the notes slide across too

        Returns:
            Status dictionary with output path, new slide number, and whether
            the layout/master were reused or copied.
        """
        if not HAS_PPTX:
            return {"error": "python-pptx/lxml not installed. Run: pip install python-pptx lxml"}

        source_resolved = resolve_office_path(source_file_path)
        target_resolved = resolve_office_path(target_file_path)
        source_path = Path(source_resolved)
        target_path = Path(target_resolved)
        if not source_path.exists():
            return {"error": f"File not found: {source_file_path}"}
        if not target_path.exists():
            return {"error": f"File not found: {target_file_path}"}

        position_text = str(position).strip().strip('"').strip("'").lower()
        if position_text not in {"end", "after"}:
            return {"error": "Invalid position. Use 'end' or 'after'."}
        if position_text == "after" and after_slide_number is None:
            return {"error": "after_slide_number is required when position='after'."}

        try:
            # Validate both packages open cleanly before doing lower-level edits.
            Presentation(str(source_path))
            target_prs = Presentation(str(target_path))
            target_slide_count_before = len(target_prs.slides)

            src_parts = _read_package(str(source_path))
            dst_parts = _read_package(str(target_path))
            src_ct_root = _parse_xml(src_parts["[Content_Types].xml"])
            dst_ct_root = _parse_xml(dst_parts["[Content_Types].xml"])

            src_slide_part, source_slide_count = _find_slide_part_by_number(src_parts, source_slide_number)
            src_layout_part = _find_slide_layout_part(src_parts, src_slide_part)

            mapping: dict[str, str] = {}
            target_layout_part, target_master_part, master_copied = _copy_master_and_layout(
                src_parts, dst_parts, src_ct_root, dst_ct_root, src_layout_part, mapping
            )
            if master_copied:
                _register_master_if_needed(dst_parts, target_master_part)

            target_slide_part = _copy_slide_part(
                src_parts,
                dst_parts,
                src_ct_root,
                dst_ct_root,
                src_slide_part,
                target_layout_part,
                mapping,
                include_notes=bool(include_notes),
            )
            new_slide_number = _insert_slide_into_presentation(
                dst_parts,
                target_slide_part,
                position=position_text,
                after_slide_number=after_slide_number,
            )

            dst_parts["[Content_Types].xml"] = _serialise_xml(dst_ct_root)

            save_path = str(output_path) if output_path else target_resolved
            _write_package(dst_parts, save_path)

            # Re-open for final validation / counts.
            final_prs = Presentation(save_path)
            return {
                "success": True,
                "source_file": source_resolved,
                "target_file": target_resolved,
                "output_file": save_path,
                "source_slide_number": source_slide_number,
                "source_slide_count": source_slide_count,
                "target_slide_count_before": target_slide_count_before,
                "target_slide_count_after": len(final_prs.slides),
                "new_slide_number": new_slide_number,
                "layout_reused": not master_copied,
                "master_copied": master_copied,
                "include_notes": bool(include_notes),
                "message": (
                    f"Imported slide {source_slide_number} from {Path(source_resolved).name} "
                    f"into {Path(save_path).name} as slide {new_slide_number}"
                ),
                "next_tools": ["pptx_list_slides", "pptx_get_slide"],
            }
        except (IndexError, ValueError) as exc:
            return {"error": str(exc)}
        except Exception as exc:  # noqa: BLE001
            return {"error": f"Failed to import slide: {type(exc).__name__}: {exc}"}
